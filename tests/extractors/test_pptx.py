"""Unit tests for :mod:`kairix.extractors.pptx` (OF-1).

Two seams are exercised:

  1. The **scripted-loader** seam — a fake ``presentation_loader``
     returning a stub presentation. Used for shape / branch tests that
     don't need the upstream library.
  2. The **real-library** seam — invokes the actual
     :mod:`pptx` library against a recorded fixture under
     ``tests/fixtures/extractors/sample.pptx``. Skipped when the
     optional ``pptx`` extra is not installed.

Sabotage-proof per test (see file ``# Sabotage:`` comments below each
assertion block):

  * ``test_extract_emits_one_page_per_slide`` — flipping
    :meth:`extract` to skip slide enumeration (or break out of the
    ``_walk_slides`` loop early) leaves ``len(doc.pages) < 3`` and
    breaks the assertion.
  * ``test_slide_titles_appear_in_page_text`` — flipping
    :func:`_slide_title` to return ``""`` for all slides leaves the
    page text without the title heading and breaks the substring check.
  * ``test_speaker_notes_appear_in_markdown`` — sabotage-proof
    executed: removing the notes-collection branch makes the assertion
    fail, confirming the test is real.
  * ``test_quality_ok_false_for_empty_deck`` — flipping
    :meth:`quality_ok` to ``return True`` breaks the assertion.
  * ``test_real_pptx_fixture_round_trips`` — flipping the temp-file
    write to drop the bytes (or pointing the loader at an empty path)
    raises before the assertion fires.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any

import pytest

from kairix.extractors import ExtractedDocument
from kairix.extractors.pptx import PptxExtractor, make_extractor, version

pytestmark = pytest.mark.unit

FIXTURE_PPTX = Path(__file__).parent.parent / "fixtures" / "extractors" / "sample.pptx"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


# ---------------------------------------------------------------------------
# Stub presentation shapes — match the wire-shape Protocol the production
# code reads off ``pptx.Presentation``.
# ---------------------------------------------------------------------------


@dataclass
class _StubShapes:
    shapes_list: list[Any] = field(default_factory=list)
    title_shape: Any = None

    def __iter__(self) -> Any:
        return iter(self.shapes_list)

    @property
    def title(self) -> Any:
        return self.title_shape


@dataclass
class _StubTextShape:
    text: str
    has_text_frame: bool = True
    shape_type: int = 17

    @property
    def text_frame(self) -> Any:
        return self


@dataclass
class _StubImageShape:
    """A shape whose ``shape_type == 13`` (MSO_SHAPE_TYPE.PICTURE)."""

    has_text_frame: bool = False
    shape_type: int = 13


@dataclass
class _StubTitleShape:
    text: str
    has_text_frame: bool = True
    shape_type: int = 14

    @property
    def text_frame(self) -> Any:
        return self


@dataclass
class _StubNotesTextFrame:
    text: str


@dataclass
class _StubNotesSlide:
    notes_text_frame: _StubNotesTextFrame


@dataclass
class _StubSlide:
    shapes: _StubShapes
    notes_slide: _StubNotesSlide | None = None


@dataclass
class _StubCoreProperties:
    title: str = "Sample Deck"
    author: str = "agent-alpha"
    created: datetime | None = None


@dataclass
class _StubPresentation:
    slides: list[_StubSlide]
    core_properties: _StubCoreProperties = field(default_factory=_StubCoreProperties)


def _build_three_slide_stub(*, with_notes: bool = True, with_image_on_slide_2: bool = False) -> _StubPresentation:
    """Build a 3-slide deck — title + body on each, optional notes."""
    slides: list[_StubSlide] = []
    for i in range(3):
        n = i + 1
        title = _StubTitleShape(text=f"Slide Title {n}")
        body = _StubTextShape(text=f"Body line for slide {n}")
        shapes_list: list[Any] = [title, body]
        if with_image_on_slide_2 and n == 2:
            shapes_list.append(_StubImageShape())
        shapes = _StubShapes(shapes_list=shapes_list, title_shape=title)
        notes_slide: _StubNotesSlide | None = None
        if with_notes:
            notes_slide = _StubNotesSlide(notes_text_frame=_StubNotesTextFrame(text=f"Speaker notes for slide {n}"))
        slides.append(_StubSlide(shapes=shapes, notes_slide=notes_slide))
    return _StubPresentation(slides=slides)


def _make_extractor(*, presentation: _StubPresentation | None = None) -> PptxExtractor:
    pres = presentation if presentation is not None else _build_three_slide_stub()
    return PptxExtractor(version=version, presentation_loader=lambda _path: pres)


# ---------------------------------------------------------------------------
# Factory + shape sanity
# ---------------------------------------------------------------------------


def test_factory_returns_pptx_instance() -> None:
    extractor = make_extractor()
    assert isinstance(extractor, PptxExtractor)
    assert extractor.name == "pptx"
    assert extractor.version == version


def test_version_module_level_non_empty() -> None:
    """F40 sanity — module-level ``version`` is a non-empty string."""
    assert isinstance(version, str)
    assert version.strip() != ""


# ---------------------------------------------------------------------------
# can_extract
# ---------------------------------------------------------------------------


def test_can_extract_claims_pptx_mime() -> None:
    extractor = _make_extractor()
    assert extractor.can_extract(_PPTX_MIME, b"") is True


def test_can_extract_claims_zip_magic_with_presentation_mime() -> None:
    extractor = _make_extractor()
    # Even if the operator passes a slightly off mime string, the
    # endswith("presentation") + ZIP magic combo claims it.
    assert extractor.can_extract("application/vnd.somevendor.presentation", b"PK\x03\x04") is True


def test_can_extract_rejects_bare_zip_without_mime_hint() -> None:
    extractor = _make_extractor()
    assert extractor.can_extract("application/octet-stream", b"PK\x03\x04") is False


def test_can_extract_rejects_text_mime() -> None:
    extractor = _make_extractor()
    assert extractor.can_extract("text/plain", b"hello") is False


# ---------------------------------------------------------------------------
# extract — per-slide pages + speaker notes + metadata
# ---------------------------------------------------------------------------


def test_extract_emits_one_page_per_slide() -> None:
    extractor = _make_extractor()
    doc = extractor.extract(b"PK\x03\x04" + b"x" * 256, _PPTX_MIME)
    assert isinstance(doc, ExtractedDocument)
    assert len(doc.pages) == 3


def test_slide_titles_appear_in_page_text() -> None:
    extractor = _make_extractor()
    doc = extractor.extract(b"PK\x03\x04" + b"x" * 256, _PPTX_MIME)
    for i, page in enumerate(doc.pages):
        # Page numbers are 1-indexed.
        assert page.page_number == i + 1
        assert f"Slide Title {i + 1}" in page.text
        assert page.text.startswith(f"## Slide {i + 1}:")


def test_speaker_notes_appear_in_markdown() -> None:
    """OF-1's headline win: notes survive into the document markdown."""
    extractor = _make_extractor()
    doc = extractor.extract(b"PK\x03\x04" + b"x" * 256, _PPTX_MIME)
    for n in range(1, 4):
        assert f"> **Speaker notes**: Speaker notes for slide {n}" in doc.markdown


def test_speaker_notes_absent_when_slide_has_no_notes() -> None:
    """A slide without ``notes_slide`` shouldn't emit a blockquote."""
    pres = _build_three_slide_stub(with_notes=False)
    extractor = _make_extractor(presentation=pres)
    doc = extractor.extract(b"PK\x03\x04" + b"x" * 256, _PPTX_MIME)
    assert "Speaker notes" not in doc.markdown


def test_extract_lifts_core_properties_into_metadata() -> None:
    extractor = _make_extractor()
    doc = extractor.extract(b"PK\x03\x04" + b"x" * 256, _PPTX_MIME)
    assert doc.metadata.title == "Sample Deck"
    assert doc.metadata.author == "agent-alpha"
    assert doc.metadata.page_count == 3


def test_extract_created_date_serialises_to_iso_when_present() -> None:
    pres = _build_three_slide_stub()
    pres.core_properties.created = datetime(2026, 5, 22, 9, 30, 0)
    extractor = _make_extractor(presentation=pres)
    doc = extractor.extract(b"PK\x03\x04" + b"x" * 16, _PPTX_MIME)
    assert doc.metadata.created_date is not None
    assert doc.metadata.created_date.startswith("2026-05-22")


def test_extract_flags_image_slide_via_has_images() -> None:
    pres = _build_three_slide_stub(with_image_on_slide_2=True)
    extractor = _make_extractor(presentation=pres)
    doc = extractor.extract(b"PK\x03\x04" + b"x" * 16, _PPTX_MIME)
    assert doc.pages[0].has_images is False
    assert doc.pages[1].has_images is True
    assert doc.pages[2].has_images is False


def test_extract_returns_zero_confidence_for_empty_bytes() -> None:
    extractor = _make_extractor()
    doc = extractor.extract(b"", _PPTX_MIME)
    assert doc.confidence == 0.0


# ---------------------------------------------------------------------------
# quality_ok
# ---------------------------------------------------------------------------


def test_quality_ok_true_for_non_empty_deck() -> None:
    extractor = _make_extractor()
    doc = extractor.extract(b"PK\x03\x04" + b"x" * 32, _PPTX_MIME)
    assert extractor.quality_ok(doc) is True


def test_quality_ok_false_for_empty_deck() -> None:
    """No slides → no pages → escalate."""
    pres = _StubPresentation(slides=[])
    extractor = _make_extractor(presentation=pres)
    doc = extractor.extract(b"PK\x03\x04" + b"x" * 1024, _PPTX_MIME)
    assert len(doc.pages) == 0
    assert extractor.quality_ok(doc) is False


def test_quality_ok_false_for_too_short_markdown() -> None:
    """One slide with no title / body / notes → too thin to keep."""
    sparse_slide = _StubSlide(shapes=_StubShapes(shapes_list=[], title_shape=None), notes_slide=None)
    pres = _StubPresentation(slides=[sparse_slide])
    extractor = _make_extractor(presentation=pres)
    doc = extractor.extract(b"PK\x03\x04" + b"x" * 16, _PPTX_MIME)
    # ``## Slide 1: (untitled)`` — under the 100-char floor.
    assert len(doc.markdown.strip()) < 100
    assert extractor.quality_ok(doc) is False


# ---------------------------------------------------------------------------
# Real-library tests — exercise the actual pptx package against a
# recorded fixture. Skipped when the optional extra is not installed.
# ---------------------------------------------------------------------------


def _pptx_available() -> bool:
    try:
        import pptx  # noqa: F401 — probe-only import; resolved at runtime
    except ImportError:
        return False
    return True


@pytest.mark.skipif(
    not _pptx_available(),
    reason="pptx extra not installed; install via 'pip install Kairix-agentic-knowledge-mgt[pptx]'",
)
def test_real_pptx_fixture_round_trips() -> None:
    raw = FIXTURE_PPTX.read_bytes()
    assert raw.startswith(b"PK\x03\x04")
    extractor = make_extractor()
    doc = extractor.extract(raw, _PPTX_MIME)
    assert isinstance(doc, ExtractedDocument)
    assert len(doc.pages) == 3
    # The fixture's slide-1 title carries known text — assert it survives.
    assert "Introduction" in doc.markdown
    # The fixture's slide-1 notes carry known text — assert they survive.
    assert "greet the audience" in doc.markdown
    assert extractor.quality_ok(doc) is True


@pytest.mark.skipif(
    not _pptx_available(),
    reason="pptx extra not installed; install via 'pip install Kairix-agentic-knowledge-mgt[pptx]'",
)
def test_real_pptx_fixture_metadata_lifted() -> None:
    raw = FIXTURE_PPTX.read_bytes()
    extractor = make_extractor()
    doc = extractor.extract(raw, _PPTX_MIME)
    assert doc.metadata.title == "Sample Deck"
    assert doc.metadata.author == "agent-alpha"
    assert doc.metadata.page_count == 3
