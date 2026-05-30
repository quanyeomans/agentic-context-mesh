#!/usr/bin/env python3
"""Generate 10 synthetic PPTX decks for the per-type fixture corpus.

ADR-028 measurement prereq. Uses python-pptx (the [pptx] extra). Slide
counts vary 3-20 per deck and content density varies bullet-heavy vs
sparse so chunker fitness measurement has a useful spread.

Slide 7 of `project-falcon-q3-architecture-decision.pptx` and slide 8
form the boundary-spanning canary pair — slide 7 sets up the decision,
slide 8 carries the conclusion. The canary suite asserts both surface.

Run from repo root:
    python3 scripts/reflib/generate_pptx_fixtures.py

Output: reference-library/per-type-fixtures/pptx/*.pptx (10 files)
"""

from __future__ import annotations

import random
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

_SCRIPT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(_SCRIPT_DIR.parent.parent))

from scripts.reflib._fixture_vocab import (  # noqa: E402 — sys.path tweak above
    AGENTS,
    PROJECTS,
    QUARTERS,
    TOPICS,
)

SEED = 28028  # ADR-028
N_FIXTURES = 10


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]  # title slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0] if bullets else ""
    for bullet in bullets[1:]:
        p = body.add_paragraph()
        p.text = bullet
        p.level = 0


def _add_sparse_slide(prs: Presentation, title: str, headline: str) -> None:
    """A sparse slide — title + a single short headline."""
    layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    txbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = txbox.text_frame
    tf.text = headline
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(36)


def _build_canary_deck(output: Path) -> None:
    """Build the deck that the boundary-spanning canary suite queries.

    Slide 7 sets up "Q3 architecture decision: two-track approach".
    Slide 8 delivers the conclusion ("we picked the per-type chunker
    plugin track"). A correct chunker keeps each slide as its own chunk;
    a regression that merges slides or splits them breaks the canary.
    """
    prs = Presentation()
    _add_title_slide(
        prs,
        "Project Falcon — Q3 Architecture Decision",
        "Two-track proposal review",
    )
    _add_bullet_slide(
        prs,
        "Agenda",
        [
            "Context: where we are after Q2",
            "Two-track proposal: per-type chunkers + eval scaffolding",
            "Decision needed today",
            "Sequencing and ownership",
        ],
    )
    _add_bullet_slide(
        prs,
        "Q2 recap",
        [
            "Single chunker still in place across every source type",
            "Benchmark reports one aggregate NDCG@10",
            "No per-type slice to point at when PPTX recall lags",
        ],
    )
    _add_bullet_slide(
        prs,
        "Track 1: per-type chunker plugins",
        [
            "Ship six structural chunkers",
            "Heading-aware for DOCX, slide-as-chunk for PPTX, row-as-chunk for XLSX",
            "Thread-aware for email, event-as-chunk for calendar",
        ],
    )
    _add_bullet_slide(
        prs,
        "Track 2: measurement surface",
        [
            "Per-source-type Recall@k slicing",
            "Boundary-spanning canary suite",
            "Chunk-size distribution telemetry",
            "Synthetic per-type fixture corpus",
        ],
    )
    _add_sparse_slide(prs, "Cost / risk", "Two-quarter window. Reversible per track.")
    # Slide 7 — sets up the decision
    _add_bullet_slide(
        prs,
        "Decision needed",
        [
            "Approve both tracks in parallel, or sequence?",
            "Sequencing risks a measurement-blind first track",
            "Parallel risks scope creep into observability work",
        ],
    )
    # Slide 8 — delivers the conclusion (canary partner of slide 7)
    _add_bullet_slide(
        prs,
        "Conclusion",
        [
            "Approved: parallel tracks with measurement scaffolding first",
            "Track 2 ships in this cycle; track 1 follows in Q4",
            "Owner: agent-alpha; review: end of cycle",
        ],
    )
    _add_sparse_slide(prs, "Next steps", "Track 2 dispatched. Track 1 spec by end of week.")
    prs.save(str(output))


def _build_generic_deck(
    output: Path, project: str, focus: str, quarter: str, slide_count: int, dense: bool, rng: random.Random
) -> None:
    """Build a non-canary deck with the given shape."""
    prs = Presentation()
    title_line = f"{project.title()} — {focus.replace('-', ' ').title()} ({quarter})"
    _add_title_slide(prs, title_line, f"Owner: {rng.choice(AGENTS)}")
    for slide_idx in range(1, slide_count):
        title = f"{rng.choice(TOPICS).title()} — slide {slide_idx}"
        if dense:
            bullets = [f"{rng.choice(AGENTS)}: {rng.choice(TOPICS)}" for _ in range(rng.randint(3, 6))]
            _add_bullet_slide(prs, title, bullets)
        else:
            _add_sparse_slide(prs, title, rng.choice(TOPICS))
    prs.save(str(output))


def generate(output_dir: Path) -> int:
    output_dir.mkdir(parents=True, exist_ok=True)
    rng = random.Random(SEED)
    written = 0
    # Canary deck first — pinned content.
    canary_path = output_dir / "project-falcon-q3-architecture-decision.pptx"
    _build_canary_deck(canary_path)
    written += 1
    # Remaining N-1 generic decks varying slide count + density
    remaining = N_FIXTURES - 1
    for i in range(remaining):
        project, focus = PROJECTS[i % len(PROJECTS)]
        quarter = QUARTERS[i % len(QUARTERS)]
        slide_count = rng.randint(3, 20)
        dense = bool(i % 2)
        path = output_dir / f"{project}-{focus}-{quarter.lower()}-deck-{i + 1:02d}.pptx"
        _build_generic_deck(path, project, focus, quarter, slide_count, dense, rng)
        written += 1
    return written


def main() -> int:
    repo_root = Path(__file__).resolve().parents[2]
    output_dir = repo_root / "reference-library" / "per-type-fixtures" / "pptx"
    n = generate(output_dir)
    print(f"pptx fixtures: wrote {n} files to {output_dir}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
