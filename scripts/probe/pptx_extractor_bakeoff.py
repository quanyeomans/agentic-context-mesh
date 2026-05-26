"""PPTX extractor A/B bakeoff — standalone diagnostic, no kairix dependency.

The v2026.5.26a1 SharePoint validation revealed that markitdown's PPTX
path is shape-naive — it dumps text per shape in document order but
loses connector edges, group hierarchy, merged-cell tables, and
SmartArt topology. This script runs the same PPTX through multiple
extractors so an operator can eyeball the differences and pick the
right tool for their corpus shape.

The script DELIBERATELY does NOT touch the kairix extractor framework.
It is a one-afternoon decision tool — install the candidate libraries,
run on 5-10 representative PPTX files, ``diff`` the outputs, decide.
Only after a decision should the winning extractor be wired into
``kairix/extractors/<name>/``.

Install once:
    pip install 'markitdown[pptx]' docling 'unstructured[pptx]'

Run on one file (all extractors):
    python scripts/probe/pptx_extractor_bakeoff.py deck.pptx

Run on a directory of decks, only docling + markitdown:
    for f in ~/decks/*.pptx; do
      python scripts/probe/pptx_extractor_bakeoff.py "$f" --only markitdown,docling
    done

Outputs land in ``./bakeoff_out/<deck-stem>.<extractor>.md`` —
``diff -u bakeoff_out/deck.markitdown.md bakeoff_out/deck.docling.md``
shows what each extractor preserves vs flattens.

Per the OSS connector-library evaluation methodology — spike first,
ADR second. This script IS the spike. The result determines whether
docling earns a permanent plugin slot, or markitdown stays and we
augment it with a python-pptx custom step for org-chart edges only,
or something else.
"""

from __future__ import annotations

import argparse
import sys
import time
import traceback
from collections.abc import Callable
from pathlib import Path


def run_markitdown(p: Path) -> str:
    """Baseline — what kairix ships today."""
    from markitdown import MarkItDown

    return MarkItDown().convert(str(p)).text_content


def run_docling(p: Path) -> str:
    """IBM docling — structure-preserving converter with typed document model."""
    from docling.document_converter import DocumentConverter

    return DocumentConverter().convert(str(p)).document.export_to_markdown()


def run_unstructured(p: Path) -> str:
    """Unstructured.io — element-typed JSON; rendered here as inspectable markdown
    with type tags + bounding boxes so the operator can see what structure
    survived. NOT meant for downstream consumption — meant for eyeballing."""
    from unstructured.partition.pptx import partition_pptx

    elements = partition_pptx(filename=str(p))
    lines = []
    for el in elements:
        bbox = getattr(el.metadata, "coordinates", None)
        page = getattr(el.metadata, "page_number", "?")
        bbox_str = ""
        if bbox is not None:
            # Coordinates are CoordinatesMetadata; render the points if present.
            pts = getattr(bbox, "points", None)
            if pts is not None:
                bbox_str = f" bbox={pts}"
        lines.append(f"[{type(el).__name__} p{page}{bbox_str}] {el.text}")
    return "\n".join(lines)


def run_python_pptx_shapes(p: Path) -> str:
    """python-pptx raw structural dump — shows what the underlying library
    actually exposes (shape types, positions, connector endpoints, table
    structure). Verbose but it's the source of truth for what's IN the file
    that markitdown/docling/unstructured then summarise."""
    from pptx import Presentation
    from pptx.util import Emu

    pres = Presentation(str(p))
    lines = []
    for slide_idx, slide in enumerate(pres.slides, 1):
        lines.append(f"## Slide {slide_idx}")
        for shape in slide.shapes:
            shape_type = type(shape).__name__
            pos = ""
            try:
                x = Emu(shape.left).inches
                y = Emu(shape.top).inches
                w = Emu(shape.width).inches
                h = Emu(shape.height).inches
                pos = f" @({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}in"
            except (AttributeError, TypeError):
                pass
            label = shape_type + pos
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(f"- {label}: {text[:200]}")
                    continue
            if shape.has_table:
                rows = len(shape.table.rows)
                cols = len(shape.table.columns)
                lines.append(f"- {label}: <table {rows}x{cols}>")
                continue
            connector_info = ""
            if hasattr(shape, "begin_connected_shape") and shape.begin_connected_shape is not None:
                connector_info = (
                    f" connector: {shape.begin_connected_shape.shape_id} → "
                    f"{shape.end_connected_shape.shape_id if shape.end_connected_shape else '?'}"
                )
            lines.append(f"- {label}{connector_info}")
    return "\n".join(lines)


EXTRACTORS: dict[str, Callable[[Path], str]] = {
    "markitdown": run_markitdown,
    "docling": run_docling,
    "unstructured": run_unstructured,
    "python_pptx_shapes": run_python_pptx_shapes,
}


def main() -> int:
    parser = argparse.ArgumentParser(
        description=__doc__,
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("pptx", type=Path, help="path to a .pptx file")
    parser.add_argument(
        "--only",
        default=",".join(EXTRACTORS),
        help=f"comma-separated extractors to run (default: all of {','.join(EXTRACTORS)})",
    )
    parser.add_argument(
        "--out-dir",
        type=Path,
        default=Path("./bakeoff_out"),
        help="where to write per-extractor markdown outputs (default: ./bakeoff_out/)",
    )
    args = parser.parse_args()

    if not args.pptx.exists():
        print(f"error: {args.pptx} does not exist", file=sys.stderr)
        return 1
    if args.pptx.suffix.lower() != ".pptx":
        print(f"error: {args.pptx} is not a .pptx file (suffix {args.pptx.suffix!r})", file=sys.stderr)
        return 1

    args.out_dir.mkdir(parents=True, exist_ok=True)
    targets = [name.strip() for name in args.only.split(",") if name.strip()]
    unknown = [n for n in targets if n not in EXTRACTORS]
    if unknown:
        print(f"error: unknown extractor(s) {unknown}; choices are {list(EXTRACTORS)}", file=sys.stderr)
        return 1

    print(f"# bakeoff for {args.pptx.name}")
    print(f"# out_dir: {args.out_dir}")
    print("# extractor       status   elapsed   chars     output_path")
    print(f"# {'-' * 80}")
    for name in targets:
        out_path = args.out_dir / f"{args.pptx.stem}.{name}.md"
        t0 = time.perf_counter()
        try:
            text = EXTRACTORS[name](args.pptx)
            elapsed = time.perf_counter() - t0
            out_path.write_text(text, encoding="utf-8")
            print(f"  {name:18s} OK       {elapsed:6.2f}s   {len(text):7d}   {out_path}")
        except ImportError as exc:
            print(f"  {name:18s} SKIP     -         -         (not installed: {exc.name})")
        except Exception as exc:
            elapsed = time.perf_counter() - t0
            print(f"  {name:18s} FAIL     {elapsed:6.2f}s   -         {type(exc).__name__}: {exc}")
            traceback.print_exc(file=sys.stderr)
    print()
    print("# next: diff -u bakeoff_out/<file>.markitdown.md bakeoff_out/<file>.docling.md")
    return 0


if __name__ == "__main__":
    sys.exit(main())
